"""PowerPoint document (.pptx) validation"""
import re
import logging
from pptx import Presentation


def validate_powerpoint_document(file_stream, rules):
    """Validate PowerPoint document against rules"""
    logging.info("Loading PowerPoint document...")
    prs = Presentation(file_stream)
    logging.info(f"PowerPoint loaded. Slides: {len(prs.slides)}")

    issues = []
    fixes_applied = []

    pptx_rules = [r for r in rules if r['doc_type'] in ['PowerPoint', 'Both', 'All']]
    ai_rules = [r for r in pptx_rules if r.get('use_ai', False)]
    hard_coded_rules = [r for r in pptx_rules if not r.get('use_ai', False)]

    logging.info(f"Hard-coded rules: {len(hard_coded_rules)} (AI rules skipped for PowerPoint)")

    # Hard-coded rules only — AI validation is skipped for PowerPoint
    # (AI is designed for prose; slide text produces too many false positives)
    for rule in hard_coded_rules:
        result = None
        if rule['rule_type'] == 'Font':
            result = _check_fonts(prs, rule)
        elif rule['rule_type'] in ['Language', 'Grammar', 'Punctuation']:
            result = _check_text(prs, rule)

        if result:
            for item in result.get('issues', []):
                if isinstance(item, dict) and 'rule_name' in item:
                    issues.append(item)
                else:
                    issues.append({
                        'rule_name': rule.get('title', 'Unknown'),
                        'rule_type': rule.get('rule_type', 'Unknown'),
                        'description': str(item),
                        'location': 'Presentation-wide',
                        'priority': rule.get('priority', 999)
                    })
            result_changes = result.get('changes', [])
            for item in result.get('fixes', []):
                if isinstance(item, dict) and 'rule_name' in item:
                    if result_changes:
                        item['changes'] = result_changes
                    fixes_applied.append(item)
                else:
                    fix_dict = {
                        'rule_name': rule.get('title', 'Unknown'),
                        'rule_type': rule.get('rule_type', 'Unknown'),
                        'found_value': 'Non-compliant value',
                        'fixed_value': str(item),
                        'location': 'Presentation-wide'
                    }
                    if result_changes:
                        fix_dict['changes'] = result_changes
                    fixes_applied.append(fix_dict)

    logging.info(f"PowerPoint validation complete. Issues: {len(issues)}, Fixes: {len(fixes_applied)}")
    return {'document': prs, 'issues': issues, 'fixes_applied': fixes_applied}


def _extract_text_refs(prs):
    """Extract all text runs with references for write-back"""
    text_refs = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            text_refs.append({
                                'run': run,
                                'text': run.text,
                                'slide': slide_idx + 1
                            })
    return text_refs


def _check_fonts(prs, rule):
    """Check and fix fonts in PowerPoint"""
    issues = []
    fixes = []
    expected_font = rule['expected_value']

    if rule['check_value'] == 'AllTextFont':
        issue_count = 0
        fix_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                current_font = run.font.name
                                if current_font is None or current_font != expected_font:
                                    issue_count += 1
                                    if rule['auto_fix']:
                                        run.font.name = expected_font
                                        fix_count += 1

        if issue_count > 0 and not rule['auto_fix']:
            issues.append({
                'rule_name': rule.get('title', 'All Text Font'),
                'rule_type': rule['rule_type'],
                'description': f"Found {issue_count} text runs with incorrect font (not {expected_font})",
                'location': 'Presentation-wide',
                'priority': rule.get('priority', 999)
            })
        if fix_count > 0:
            fixes.append({
                'rule_name': rule.get('title', 'All Text Font'),
                'rule_type': rule['rule_type'],
                'found_value': f'{issue_count} runs with wrong font',
                'fixed_value': expected_font,
                'location': f'Presentation-wide ({fix_count} runs)'
            })

    return {'issues': issues, 'fixes': fixes}


def _check_text(prs, rule):
    """Check and fix text issues in PowerPoint (spelling, contractions, symbols, numbers)"""
    issues = []
    fixes = []
    changes = []
    check_value = rule['check_value']
    issue_count = 0
    fix_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text or not run.text.strip():
                        continue
                    text = run.text
                    location = f'Slide {slide_idx + 1}'

                    if check_value.startswith('BritishSpelling_'):
                        american_word = check_value.replace('BritishSpelling_', '')
                        british_word = rule['expected_value']
                        pattern = r'\b' + re.escape(american_word) + r'\b'
                        matches = re.findall(pattern, text, re.IGNORECASE)
                        if matches:
                            issue_count += len(matches)
                            if rule['auto_fix']:
                                def replace_preserve_case(match, replacement=british_word):
                                    word = match.group(0)
                                    if word.isupper():
                                        return replacement.upper()
                                    elif word[0].isupper():
                                        return replacement.capitalize()
                                    return replacement
                                before = run.text
                                run.text = re.sub(pattern, replace_preserve_case, text, flags=re.IGNORECASE)
                                changes.append({'before': before, 'after': run.text, 'location': location})
                                fix_count += len(matches)

                    elif check_value.startswith('NoContraction_'):
                        contraction = check_value.replace('NoContraction_', '')
                        expanded = rule['expected_value']
                        if contraction in text:
                            count = text.count(contraction)
                            issue_count += count
                            if rule['auto_fix']:
                                before = run.text
                                run.text = text.replace(contraction, expanded)
                                changes.append({'before': before, 'after': run.text, 'location': location})
                                fix_count += count

                    elif check_value == 'NoAmpersand':
                        if '&' in text:
                            count = text.count('&')
                            issue_count += count
                            if rule['auto_fix']:
                                before = run.text
                                run.text = text.replace('&', 'and')
                                changes.append({'before': before, 'after': run.text, 'location': location})
                                fix_count += count

                    elif check_value == 'PercentSymbol':
                        percent_matches = re.findall(r'\d+%', text)
                        if percent_matches:
                            issue_count += len(percent_matches)
                            if rule['auto_fix']:
                                before = run.text
                                run.text = re.sub(r'(\d+)%', r'\1 percent', text)
                                changes.append({'before': before, 'after': run.text, 'location': location})
                                fix_count += len(percent_matches)

                    elif check_value == 'NoApostrophePlurals':
                        apos_matches = re.findall(r"\b[A-Z]{2,}'s\b", text)
                        if apos_matches:
                            issue_count += len(apos_matches)

                    elif check_value == 'NumberCommas':
                        num_matches = re.findall(r'\b\d{4,}\b', text)
                        num_matches = [m for m in num_matches if not (1900 <= int(m) <= 2099)]
                        if num_matches:
                            issue_count += len(num_matches)
                            if rule['auto_fix']:
                                before = run.text
                                for match in num_matches:
                                    formatted = '{:,}'.format(int(match))
                                    run.text = run.text.replace(match, formatted)
                                    fix_count += 1
                                changes.append({'before': before, 'after': run.text, 'location': location})

                    elif check_value == 'Word_toward':
                        toward_matches = re.findall(r'\btowards\b', text, re.IGNORECASE)
                        if toward_matches:
                            issue_count += len(toward_matches)
                            if rule['auto_fix']:
                                before = run.text
                                run.text = re.sub(r'\btowards\b', 'toward', text, flags=re.IGNORECASE)
                                changes.append({'before': before, 'after': run.text, 'location': location})
                                fix_count += len(toward_matches)

                    elif check_value == 'AvoidEtc':
                        etc_matches = re.findall(r'\betc\.?\b', text, re.IGNORECASE)
                        if etc_matches:
                            issue_count += len(etc_matches)

    label = rule.get('title', check_value)
    if issue_count > 0:
        issues.append(f"Found {issue_count} instances of '{label}' violations")
    if fix_count > 0:
        fixes.append(f"Fixed {fix_count} instances for '{label}'")

    return {'issues': issues, 'fixes': fixes, 'changes': changes}
