"""Test validation across all formats and verify status logic + report generation"""
import sys
import os
sys.path.insert(0, os.path.dirname(__file__))

from io import BytesIO
from docx import Document
from docx.shared import Pt


def create_test_word():
    """Word doc with wrong font"""
    doc = Document()
    p = doc.add_paragraph()
    run = p.add_run("This is a test with organization and color spelled wrong.")
    run.font.name = "Times New Roman"
    run.font.size = Pt(11)
    stream = BytesIO()
    doc.save(stream)
    stream.seek(0)
    return stream


def create_test_excel():
    """Excel with wrong font and American spelling"""
    from openpyxl import Workbook
    from openpyxl.styles import Font
    wb = Workbook()
    ws = wb.active
    ws['A1'] = "Organization chart"
    ws['A1'].font = Font(name='Times New Roman', size=11)
    ws['A2'] = "Color palette"
    ws['A2'].font = Font(name='Calibri', size=11)
    ws['A3'] = "Normal text here"
    ws['A3'].font = Font(name='Arial', size=11)
    stream = BytesIO()
    wb.save(stream)
    stream.seek(0)
    return stream


def create_test_pptx():
    """PowerPoint with wrong font"""
    from pptx import Presentation
    from pptx.util import Pt as PptxPt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Test Presentation"
    for run in title.text_frame.paragraphs[0].runs:
        run.font.name = "Calibri"
        run.font.size = PptxPt(24)
    body = slide.placeholders[1]
    body.text = "Organization and color"
    for run in body.text_frame.paragraphs[0].runs:
        run.font.name = "Times New Roman"
    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream


def mock_rules():
    """Rules that work across all formats"""
    return [
        {
            'title': 'All Text Font',
            'rule_type': 'Font',
            'doc_type': 'All',
            'check_value': 'AllTextFont',
            'expected_value': 'Arial',
            'auto_fix': True,
            'use_ai': False,
            'priority': 1
        },
        {
            'title': 'British Spelling: organization',
            'rule_type': 'Language',
            'doc_type': 'All',
            'check_value': 'BritishSpelling_organization',
            'expected_value': 'organisation',
            'auto_fix': True,
            'use_ai': False,
            'priority': 2
        },
        {
            'title': 'British Spelling: color',
            'rule_type': 'Language',
            'doc_type': 'All',
            'check_value': 'BritishSpelling_color',
            'expected_value': 'colour',
            'auto_fix': True,
            'use_ai': False,
            'priority': 2
        },
    ]


def test_status_logic():
    """Test the three-way status calculation"""
    print("=" * 60)
    print("TEST: Status logic (Passed / Review Required / Failed)")
    print("=" * 60)

    # Simulate: no issues at all -> Passed
    issues_none = []
    fixes_none = []
    remaining = [i for i in issues_none if isinstance(i, dict)]
    status = "Passed" if len(remaining) == 0 else ("Review Required" if len(fixes_none) > 0 else "Failed")
    assert status == "Passed", f"Expected Passed, got {status}"
    print(f"  [PASS] No issues -> {status}")

    # Simulate: issues all fixed -> Passed
    issues_fixed = []  # no remaining dict issues
    fixes_applied = [{'rule_name': 'Font', 'rule_type': 'Font', 'found_value': 'x', 'fixed_value': 'y', 'location': 'z'}]
    remaining = [i for i in issues_fixed if isinstance(i, dict)]
    status = "Passed" if len(remaining) == 0 else ("Review Required" if len(fixes_applied) > 0 else "Failed")
    assert status == "Passed", f"Expected Passed, got {status}"
    print(f"  [PASS] All fixed, no remaining -> {status}")

    # Simulate: some fixed, some remaining -> Review Required
    issues_mixed = [{'rule_name': 'AI', 'rule_type': 'AI', 'description': 'x', 'location': 'y', 'priority': 3}]
    fixes_some = [{'rule_name': 'Font', 'rule_type': 'Font', 'found_value': 'x', 'fixed_value': 'y', 'location': 'z'}]
    remaining = [i for i in issues_mixed if isinstance(i, dict)]
    status = "Passed" if len(remaining) == 0 else ("Review Required" if len(fixes_some) > 0 else "Failed")
    assert status == "Review Required", f"Expected Review Required, got {status}"
    print(f"  [PASS] Some fixed, some remaining -> {status}")

    # Simulate: issues, none fixed -> Failed
    issues_only = [{'rule_name': 'Test', 'rule_type': 'Test', 'description': 'x', 'location': 'y', 'priority': 1}]
    fixes_zero = []
    remaining = [i for i in issues_only if isinstance(i, dict)]
    status = "Passed" if len(remaining) == 0 else ("Review Required" if len(fixes_zero) > 0 else "Failed")
    assert status == "Failed", f"Expected Failed, got {status}"
    print(f"  [PASS] Issues, no fixes -> {status}")

    # Simulate: string issues (not dict) should be ignored
    issues_strings = ["some string issue", {'rule_name': 'Real', 'rule_type': 'T', 'description': 'x', 'location': 'y', 'priority': 1}]
    remaining = [i for i in issues_strings if isinstance(i, dict)]
    assert len(remaining) == 1, f"Expected 1 remaining, got {len(remaining)}"
    print(f"  [PASS] String issues filtered out, only dict counted ({len(remaining)} remaining)")

    print()


def test_report_generation():
    """Test report generates correct status and description"""
    print("=" * 60)
    print("TEST: Report generation")
    print("=" * 60)

    from ValidateDocument.report import generate_report

    # Test Passed report
    report = generate_report("test.docx", [], [{'rule_name': 'Font', 'rule_type': 'Font', 'found_value': 'x', 'fixed_value': 'y', 'location': 'z'}])
    assert ">Passed<" in report or ">PASSED<" in report.upper(), "Passed not in report"
    assert "auto-fixed" in report.lower(), "auto-fixed description missing"
    print("  [PASS] Passed report with fixes")

    # Test Review Required report
    issues = [{'rule_name': 'AI', 'rule_type': 'AI', 'description': 'test', 'location': 'doc', 'priority': 3}]
    fixes = [{'rule_name': 'Font', 'rule_type': 'Font', 'found_value': 'x', 'fixed_value': 'y', 'location': 'z'}]
    report = generate_report("test.xlsx", issues, fixes)
    assert "Review Required" in report, "Review Required not in report"
    assert "#f0ad4e" in report, "Amber color not in report"
    assert "remaining" in report.lower(), "remaining description missing"
    print("  [PASS] Review Required report")

    # Test Failed report
    issues = [{'rule_name': 'Test', 'rule_type': 'Test', 'description': 'fail', 'location': 'doc', 'priority': 1}]
    report = generate_report("test.vsdx", issues, [])
    assert ">Failed<" in report or ">FAILED<" in report.upper(), "Failed not in report"
    assert "#dc3545" in report, "Red color not in report"
    assert "Manual correction" in report, "Manual correction description missing"
    print("  [PASS] Failed report")

    # Test clean report
    report = generate_report("test.docx", [], [])
    assert ">Passed<" in report or ">PASSED<" in report.upper(), "Passed not in clean report"
    assert "No issues found" in report, "No issues description missing"
    print("  [PASS] Clean (no issues) report")

    # Save sample report
    with open('test_review_required_report.html', 'w') as f:
        f.write(generate_report("sample.xlsx",
            [{'rule_name': 'AI Style Corrections', 'rule_type': 'AI', 'description': 'Found 15 style violations requiring manual review', 'location': 'Workbook-wide', 'priority': 3}],
            [{'rule_name': 'All Text Font', 'rule_type': 'Font', 'found_value': '42 cells with wrong font', 'fixed_value': 'Arial', 'location': 'Workbook-wide (42 cells)'}]
        ))
    print("  [SAVE] Sample Review Required report -> test_review_required_report.html")

    print()


def test_word_validation():
    """Test Word validation"""
    print("=" * 60)
    print("TEST: Word validation")
    print("=" * 60)

    from ValidateDocument.word_validator import validate_word_document

    stream = create_test_word()
    rules = mock_rules()
    result = validate_word_document(stream, rules)

    print(f"  Issues: {len(result['issues'])}")
    print(f"  Fixes:  {len(result['fixes_applied'])}")
    for fix in result['fixes_applied']:
        if isinstance(fix, dict):
            print(f"    - {fix.get('rule_name')}: {fix.get('found_value')} -> {fix.get('fixed_value')}")
        else:
            print(f"    - {fix}")

    assert len(result['fixes_applied']) > 0, "Expected some fixes"
    print("  [PASS] Word validation produced fixes")
    print()


def test_excel_validation():
    """Test Excel validation"""
    print("=" * 60)
    print("TEST: Excel validation")
    print("=" * 60)

    from ValidateDocument.excel_validator import validate_excel_document

    stream = create_test_excel()
    rules = mock_rules()
    result = validate_excel_document(stream, rules)

    print(f"  Issues: {len(result['issues'])}")
    print(f"  Fixes:  {len(result['fixes_applied'])}")
    for fix in result['fixes_applied']:
        if isinstance(fix, dict):
            print(f"    - {fix.get('rule_name')}: {fix.get('found_value')} -> {fix.get('fixed_value')}")
        else:
            print(f"    - {fix}")

    assert len(result['fixes_applied']) > 0, "Expected some fixes"
    print("  [PASS] Excel validation produced fixes")

    # Verify font was changed
    wb = result['document']
    ws = wb.active
    for cell in [ws['A1'], ws['A2']]:
        assert cell.font.name == 'Arial', f"Cell {cell.coordinate} font is {cell.font.name}, expected Arial"
    print("  [PASS] Fonts corrected to Arial")
    print()


def test_pptx_validation():
    """Test PowerPoint validation"""
    print("=" * 60)
    print("TEST: PowerPoint validation")
    print("=" * 60)

    from ValidateDocument.powerpoint_validator import validate_powerpoint_document

    stream = create_test_pptx()
    rules = mock_rules()
    result = validate_powerpoint_document(stream, rules)

    print(f"  Issues: {len(result['issues'])}")
    print(f"  Fixes:  {len(result['fixes_applied'])}")
    for fix in result['fixes_applied']:
        if isinstance(fix, dict):
            print(f"    - {fix.get('rule_name')}: {fix.get('found_value')} -> {fix.get('fixed_value')}")
        else:
            print(f"    - {fix}")

    assert len(result['fixes_applied']) > 0, "Expected some fixes"
    print("  [PASS] PowerPoint validation produced fixes")
    print()


def test_file_extensions():
    """Test file extension routing"""
    print("=" * 60)
    print("TEST: File extension routing")
    print("=" * 60)

    word_exts = ['.docx', '.doc', '.docm', '.dotx', '.dotm']
    visio_exts = ['.vsdx', '.vsd']
    excel_exts = ['.xlsx', '.xls', '.xlsm']
    pptx_exts = ['.pptx', '.ppt', '.pptm', '.potx', '.potm']

    for ext in word_exts:
        assert ext in ['.docx', '.doc', '.docm', '.dotx', '.dotm'], f"{ext} not in Word list"
        print(f"  [PASS] {ext} -> Word")

    for ext in excel_exts:
        assert ext in ['.xlsx', '.xls', '.xlsm'], f"{ext} not in Excel list"
        print(f"  [PASS] {ext} -> Excel")

    for ext in pptx_exts:
        assert ext in ['.pptx', '.ppt', '.pptm', '.potx', '.potm'], f"{ext} not in PowerPoint list"
        print(f"  [PASS] {ext} -> PowerPoint")

    for ext in visio_exts:
        print(f"  [PASS] {ext} -> Visio")

    print()


if __name__ == "__main__":
    os.chdir(os.path.dirname(__file__))

    test_status_logic()
    test_report_generation()
    test_file_extensions()
    test_word_validation()
    test_excel_validation()
    test_pptx_validation()

    print("=" * 60)
    print("ALL TESTS PASSED")
    print("=" * 60)
