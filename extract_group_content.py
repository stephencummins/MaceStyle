"""
Extract content from grouped shapes in PowerPoint
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

def extract_group_shapes(shape, level=0):
    """Recursively extract content from grouped shapes"""
    indent = "  " * level
    results = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        results.append(f"{indent}[GROUP: {shape.name}]")
        for sub_shape in shape.shapes:
            results.extend(extract_group_shapes(sub_shape, level + 1))
    else:
        info = f"{indent}- {shape.shape_type.name}"
        if hasattr(shape, 'text') and shape.text:
            info += f": {shape.text.strip()}"
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            from pptx.util import Inches
            info += f" (Size: {shape.width/Inches(1):.2f}\" × {shape.height/Inches(1):.2f}\")"
        results.append(info)

        # Check for tables in the shape
        if shape.has_table:
            results.append(f"{indent}  [TABLE with {len(shape.table.rows)} rows × {len(shape.table.columns)} columns]")
            for row_idx, row in enumerate(shape.table.rows):
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                if any(row_text):
                    results.append(f"{indent}    Row {row_idx + 1}: {' | '.join(row_text)}")

    return results

def main():
    pptx_path = "/Users/stephen/Downloads/Process Map_PowerPoint Template.pptx"

    print("=" * 70)
    print("EXTRACTING GROUPED CONTENT FROM SLIDE 2 (GUIDANCE)")
    print("=" * 70)
    print()

    try:
        prs = Presentation(pptx_path)

        if len(prs.slides) < 2:
            print("Error: Not enough slides in presentation")
            return False

        slide = prs.slides[1]  # Slide 2 (0-indexed)

        print("Analyzing all shapes on Slide 2:\n")

        all_content = []

        for idx, shape in enumerate(slide.shapes, 1):
            print(f"\n{'=' * 70}")
            print(f"Shape {idx}: {shape.name}")
            print('=' * 70)

            content = extract_group_shapes(shape)
            for line in content:
                print(line)
                all_content.append(line)

        # Save to file
        output_path = "/Users/stephen/Projects/MaceStyle/slide2_guidance_content.txt"
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("SLIDE 2 - GUIDANCE PAGE CONTENT\n")
            f.write("=" * 70 + "\n\n")
            for line in all_content:
                f.write(line + "\n")

        print(f"\n\n✓ Content saved to: {output_path}")
        print()

        return True

    except Exception as e:
        print(f"\n✗ Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
