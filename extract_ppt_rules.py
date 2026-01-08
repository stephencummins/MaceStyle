"""
Extract rules from PowerPoint template
"""
from pptx import Presentation
import sys

def extract_all_text(pptx_path):
    """Extract all text from PowerPoint file"""

    print("=" * 70)
    print("EXTRACTING TEXT FROM POWERPOINT")
    print("=" * 70)
    print()

    try:
        prs = Presentation(pptx_path)

        print(f"File: {pptx_path}")
        print(f"Total slides: {len(prs.slides)}")
        print()

        all_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"\n{'=' * 70}")
            print(f"SLIDE {slide_num}")
            print('=' * 70)

            slide_text = []

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                        print(f"\n{text}")

                # Check for tables
                if shape.has_table:
                    print("\n[TABLE CONTENT:]")
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            print(" | ".join(row_text))

            all_text.extend(slide_text)

        print()
        print("=" * 70)
        print("EXTRACTION COMPLETE")
        print("=" * 70)
        print(f"Total text blocks extracted: {len(all_text)}")
        print()

        # Save to file
        output_path = "/Users/stephen/Projects/MaceStyle/ppt_extracted_text.txt"
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(f"Extracted from: {pptx_path}\n")
            f.write(f"Total slides: {len(prs.slides)}\n")
            f.write("=" * 70 + "\n\n")

            for i, slide in enumerate(prs.slides, 1):
                f.write(f"\n{'=' * 70}\n")
                f.write(f"SLIDE {i}\n")
                f.write('=' * 70 + '\n\n')

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            f.write(f"{text}\n\n")

                    if shape.has_table:
                        f.write("\n[TABLE CONTENT:]\n")
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                f.write(" | ".join(row_text) + "\n")
                        f.write("\n")

        print(f"✓ Full extraction saved to: {output_path}")
        print()

        return True

    except Exception as e:
        print(f"\n✗ Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    pptx_path = "/Users/stephen/Downloads/Process Map_PowerPoint Template.pptx"
    success = extract_all_text(pptx_path)
    sys.exit(0 if success else 1)
