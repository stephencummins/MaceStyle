"""
Detailed PowerPoint analysis to extract formatting rules and structure
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import sys

def analyze_ppt_structure(pptx_path):
    """Analyze PowerPoint structure and extract formatting patterns"""

    print("=" * 70)
    print("DETAILED POWERPOINT ANALYSIS")
    print("=" * 70)
    print()

    try:
        prs = Presentation(pptx_path)

        print(f"File: {pptx_path}")
        print(f"Slide dimensions: {prs.slide_width/Inches(1):.2f}\" x {prs.slide_height/Inches(1):.2f}\"")
        print(f"Total slides: {len(prs.slides)}")
        print()

        # Analyze slide 2 (guidance page) in detail
        if len(prs.slides) >= 2:
            print("=" * 70)
            print("SLIDE 2 (GUIDANCE PAGE) - DETAILED ANALYSIS")
            print("=" * 70)

            slide = prs.slides[1]  # Slide 2 (0-indexed)

            print(f"\nTotal shapes on this slide: {len(slide.shapes)}")
            print()

            for idx, shape in enumerate(slide.shapes, 1):
                print(f"\nShape {idx}:")
                print(f"  Type: {shape.shape_type}")
                print(f"  Name: {shape.name if hasattr(shape, 'name') else 'N/A'}")

                # Position and size
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    print(f"  Position: ({shape.left/Inches(1):.2f}\", {shape.top/Inches(1):.2f}\")")
                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    print(f"  Size: {shape.width/Inches(1):.2f}\" x {shape.height/Inches(1):.2f}\"")

                # Text content
                if hasattr(shape, "text") and shape.text:
                    print(f"  Text: {shape.text[:100]}")  # First 100 chars

                # Notes
                if hasattr(shape, "text_frame"):
                    print(f"  Has text frame: Yes")

                # Tables
                if shape.has_table:
                    print(f"  Has table: Yes")
                    print(f"  Table size: {len(shape.table.rows)} rows x {len(shape.table.columns)} cols")

                # Images/Pictures
                if hasattr(shape, "image"):
                    print(f"  Contains image: Yes")

        # Check for slide notes
        print("\n" + "=" * 70)
        print("CHECKING FOR SLIDE NOTES")
        print("=" * 70)

        for i, slide in enumerate(prs.slides, 1):
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame'):
                    notes_text = notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        print(f"\nSlide {i} Notes:")
                        print(notes_text)

        # Analyze swim lane structure from process map slides
        print("\n" + "=" * 70)
        print("ANALYZING PROCESS MAP STRUCTURE")
        print("=" * 70)

        swim_lanes = set()
        for slide in prs.slides[3:]:  # Process map slides start at slide 4
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and len(text) < 100:  # Swim lane labels are short
                        swim_lanes.add(text)

        print("\nIdentified swim lanes:")
        for lane in sorted(swim_lanes):
            if lane and lane not in ["Document Ref:", "Insert Process Title"]:
                print(f"  - {lane}")

        # Look for shape patterns
        print("\n" + "=" * 70)
        print("SHAPE PATTERNS (from slide 3 - template shapes)")
        print("=" * 70)

        if len(prs.slides) >= 3:
            template_slide = prs.slides[2]  # Slide 3
            print(f"\nTemplate shapes available:")

            for shape in template_slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and text not in ["Copy and paste these shapes onto your process map in the appropriate swim lane."]:
                        if hasattr(shape, 'width') and hasattr(shape, 'height'):
                            print(f"  - {text}: {shape.width/Inches(1):.2f}\" x {shape.height/Inches(1):.2f}\"")

        # Generate inferred rules
        print("\n" + "=" * 70)
        print("INFERRED FORMATTING RULES")
        print("=" * 70)
        print()
        print("Based on the template structure, potential validation rules:")
        print()
        print("1. SLIDE SIZE:")
        print(f"   - Standard: {prs.slide_width/Inches(1):.2f}\" x {prs.slide_height/Inches(1):.2f}\"")
        print()
        print("2. SWIM LANE STRUCTURE:")
        print("   - New Hospitals Programme")
        print("   - NHS")
        print("   - Healthy Delivery Partnership")
        print("   - Delivery Team")
        print("   - Contractor/Supply Chain")
        print()
        print("3. REQUIRED FIELDS:")
        print("   - Document Ref: (must be present)")
        print("   - Process Title (must be present)")
        print("   - Function - Sub Function label")
        print()
        print("4. PROCESS SHAPES:")
        print("   - Start shape (specific style)")
        print("   - Set Up shape (specific style)")
        print("   - Approve shape (specific style)")
        print("   - Process boxes with XX-000001 format reference")
        print()

        # Save detailed analysis
        output_path = "/Users/stephen/Projects/MaceStyle/ppt_detailed_analysis.txt"
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("PROCESS MAP TEMPLATE ANALYSIS\n")
            f.write("=" * 70 + "\n\n")
            f.write(f"Template: {pptx_path}\n")
            f.write(f"Slide Size: {prs.slide_width/Inches(1):.2f}\" x {prs.slide_height/Inches(1):.2f}\"\n\n")

            f.write("\nREQUIRED SWIM LANES:\n")
            f.write("  1. New Hospitals Programme\n")
            f.write("  2. NHS\n")
            f.write("  3. Healthy Delivery Partnership\n")
            f.write("  4. Delivery Team\n")
            f.write("  5. Contractor/Supply Chain\n\n")

            f.write("\nFUNCTIONS:\n")
            functions = ["PMO", "People", "Digital", "Commercial", "Technical Services",
                        "Delivery", "Industrialisation", "PSMO", "Operations (Transformation)"]
            for i, func in enumerate(functions, 1):
                f.write(f"  {i}. {func}\n")

        print(f"\n✓ Detailed analysis saved to: {output_path}")
        print()

        return True

    except Exception as e:
        print(f"\n✗ Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    pptx_path = "/Users/stephen/Downloads/Process Map_PowerPoint Template.pptx"
    success = analyze_ppt_structure(pptx_path)
    sys.exit(0 if success else 1)
