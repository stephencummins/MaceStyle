"""Generate the Mace Style Validator executive briefing deck (MaceWay Control Centre branding)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "maceway-logo.png")
OUT = os.path.join(HERE, "MaceStyleValidator-Exec-Briefing.pptx")

# Brand palette
GREY  = RGBColor(0x79, 0x75, 0x74)   # MaceWay Control Centre grey
DARK  = RGBColor(0x40, 0x40, 0x40)
MUTED = RGBColor(0x6C, 0x75, 0x7D)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x52)
TEXT  = RGBColor(0x33, 0x33, 0x33)
FONT  = "Arial"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False):
    run.font.name = FONT; run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def add_logo(slide, left, top, height):
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, left, top, height=height)


def title_bar(slide, title, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREY; bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.55)
    p = tf.paragraphs[0]
    if kicker:
        rk = p.add_run(); rk.text = kicker.upper() + "\n"; _set(rk, 11, RGBColor(0xDD, 0xDD, 0xDD), bold=True)
    r = p.add_run(); r.text = title; _set(r, 26, WHITE, bold=True)
    add_logo(slide, EMU_W - Inches(2.05), Inches(0.34), Inches(0.48))


def bullets(slide, items, left=Inches(0.7), top=Inches(1.5),
            width=Inches(12.0), height=Inches(5.4), size=18, gap=8):
    tf = textbox(slide, left, top, width, height)
    first = True
    for it in items:
        text, level = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(gap); p.level = max(level, 0)
        if level == -1:
            r = p.add_run(); r.text = text; _set(r, size + 2, GREY, bold=True); p.space_before = Pt(10)
        else:
            prefix = "•  " if level == 0 else "–  "
            r = p.add_run(); r.text = prefix + text
            _set(r, size - (level * 2), TEXT if level == 0 else MUTED)
    return tf


def card(slide, left, top, width, height, heading, body, accent=GREY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = accent
    box.line.width = Pt(1.25); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = heading; _set(r, 16, accent, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(6); r2 = p2.add_run(); r2.text = body; _set(r2, 13, TEXT)


def footer(slide):
    page = len(prs.slides._sldIdLst)
    tf = textbox(slide, Inches(0.5), Inches(7.05), Inches(11.6), Inches(0.35))
    r = tf.paragraphs[0].add_run()
    r.text = "Mace Style Validator  ·  Executive briefing  ·  Confidential"; _set(r, 9, MUTED)
    pg = textbox(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    pp = pg.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run(); rr.text = str(page); _set(rr, 9, MUTED)


# ---- diagram primitives ---------------------------------------------------
def node(slide, x, y, w, h, title, sub="", fill=GREY, txt=WHITE, tsz=13, ssz=10):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill; b.line.fill.background(); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = title; _set(r, tsz, txt, bold=True)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; r2 = p2.add_run(); r2.text = sub; _set(r2, ssz, txt)
    return b


def container(slide, x, y, w, h, label, border):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.background(); box.line.color.rgb = border; box.line.width = Pt(2.25); box.shadow.inherit = False
    tab = textbox(slide, x + Inches(0.18), y + Inches(0.08), w - Inches(0.3), Inches(0.34))
    r = tab.paragraphs[0].add_run(); r.text = label; _set(r, 12, border, bold=True)
    return box


def arrow(slide, shape, x, y, w, h, color=MUTED):
    a = slide.shapes.add_shape(shape, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    return a


def note_strip(slide, x, y, w, h, text, border=GREEN):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = border
    box.line.width = Pt(1.25); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    r = tf.paragraphs[0].add_run(); r.text = text; _set(r, 13, TEXT, bold=True)


# ================================================================ 1. TITLE
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
bg.fill.solid(); bg.fill.fore_color.rgb = GREY; bg.line.fill.background(); bg.shadow.inherit = False
add_logo(s, Inches(0.9), Inches(1.1), Inches(1.5))
tf = textbox(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(2.8))
r = tf.paragraphs[0].add_run(); r.text = "Mace Style Validator"; _set(r, 46, WHITE, bold=True)
p = tf.add_paragraph(); r = p.add_run()
r.text = "Automated document compliance for the Mace Way Control Centre"; _set(r, 22, RGBColor(0xEC, 0xEC, 0xEC))
p = tf.add_paragraph(); p.space_before = Pt(22); r = p.add_run()
r.text = "Executive briefing for Tarek ElHadidi, Chief Digital & Information Officer"; _set(r, 16, WHITE, bold=True)
p = tf.add_paragraph(); r = p.add_run()
r.text = "Prepared by Stephen Cummins  ·  June 2026"; _set(r, 13, RGBColor(0xDD, 0xDD, 0xDD))

# ================================================================ 2. OPPORTUNITY
s = prs.slides.add_slide(BLANK)
title_bar(s, "The opportunity", "Why this matters")
bullets(s, [
    "The Mace Way Control Centre is the single source of truth for how Mace delivers — hundreds of governance documents that must follow the Control Centre Writing Style Guide.",
    "Checking style by hand is slow, inconsistent, and doesn't scale across Word, Excel, PowerPoint and Visio.",
    "Inconsistent documents quietly erode the authority and professionalism of the Control Centre.",
    "There was no automated, auditable way to enforce the standard at the point documents are created.",
], gap=14)
footer(s)

# ================================================================ 3. SOLUTION
s = prs.slides.add_slide(BLANK)
title_bar(s, "What we built", "The solution")
tf = textbox(s, Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.9))
r = tf.paragraphs[0].add_run()
r.text = ("An automated style validator, built into SharePoint, that checks and "
          "auto-fixes documents against the Mace Way Writing Style Guide — in one click.")
_set(r, 18, TEXT, bold=True)
w, h, g, x0, y0 = Inches(3.95), Inches(1.5), Inches(0.25), Inches(0.7), Inches(2.55)
for i, (hd, bd) in enumerate([
    ("One click", "Run it straight from the document library — no new tools to learn."),
    ("Auto-fix", "Common issues are corrected automatically in a clean copy."),
    ("Branded report", "Every run produces a Mace-branded report, linked on the document."),
    ("Full record", "Pass / review / fail, with issues found and fixed, logged for audit."),
    ("Multi-format", "Word, Excel, PowerPoint and Visio — one engine."),
    ("Business-owned rules", "199 rules today, editable in SharePoint with no code change."),
]):
    card(s, x0 + (i % 3) * (w + g), y0 + (i // 3) * (h + Inches(0.25)), w, h, hd, bd)
footer(s)

# ================================================================ 4. HOW USERS RUN IT
s = prs.slides.add_slide(BLANK)
title_bar(s, "How a user runs it", "User guide")
bullets(s, [
    'In the document library, set a document\'s status to "Validate Now".',
    "The validator checks it against the live style rules held in SharePoint.",
    "Common issues are auto-fixed in a corrected copy of the document.",
    "A Mace-branded report is generated and linked directly on the document.",
    "The outcome — pass, review-required or fail, with counts — is recorded in the Validation Results list.",
    "Reports and results carry links back to the source document and library.",
], gap=12)
footer(s)

# ================================================================ 5. PROCESS FLOW (diagram)
s = prs.slides.add_slide(BLANK)
title_bar(s, "From upload to compliant — automatically", "How it works")
steps = [
    ('1. "Validate Now"', "User flags the document\nin the library"),
    ("2. Flow triggers", "Sends the document to\nthe validation engine"),
    ("3. Checked", "Validated against the\nlive style rules"),
    ("4. Auto-fixed", "Issues corrected in a\nclean copy"),
    ("5. Report", "Branded report generated\nand linked on the doc"),
    ("6. Recorded", "Pass / review / fail\nlogged for audit"),
]
nw, nh = Inches(3.7), Inches(1.4)
top_y, bot_y = Inches(1.75), Inches(3.95)
xs = [Inches(0.6), Inches(4.8), Inches(9.0)]
fills = [GREY, GREY, GREY, GREEN, GREEN, GREEN]
# top row L->R (0,1,2)
for i in range(3):
    node(s, xs[i], top_y, nw, nh, steps[i][0], steps[i][1], fill=fills[i], tsz=15, ssz=11)
arrow(s, MSO_SHAPE.RIGHT_ARROW, Inches(4.3), Inches(2.32), Inches(0.5), Inches(0.36))
arrow(s, MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(2.32), Inches(0.5), Inches(0.36))
# down from node3 to node4 (both at right column)
arrow(s, MSO_SHAPE.DOWN_ARROW, Inches(10.65), Inches(3.25), Inches(0.4), Inches(0.6))
# bottom row R->L (3 at right, 4 mid, 5 left)
node(s, xs[2], bot_y, nw, nh, steps[3][0], steps[3][1], fill=fills[3], tsz=15, ssz=11)
node(s, xs[1], bot_y, nw, nh, steps[4][0], steps[4][1], fill=fills[4], tsz=15, ssz=11)
node(s, xs[0], bot_y, nw, nh, steps[5][0], steps[5][1], fill=fills[5], tsz=15, ssz=11)
arrow(s, MSO_SHAPE.LEFT_ARROW, Inches(8.5), Inches(4.52), Inches(0.5), Inches(0.36))
arrow(s, MSO_SHAPE.LEFT_ARROW, Inches(4.3), Inches(4.52), Inches(0.5), Inches(0.36))
note_strip(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(0.8),
           "Steps 2–6 are fully automated. The user only sets “Validate Now” — everything else happens in seconds.", border=GREEN)
footer(s)

# ================================================================ 6. WHAT IT CHECKS
s = prs.slides.add_slide(BLANK)
title_bar(s, "What it checks", "Coverage")
bullets(s, [
    ("Style rules enforced", -1),
    ("British English spelling (colour, analyse, centre …)", 0),
    ("Contraction expansion, symbol replacement (& → and, % → percent)", 0),
    ("Number formatting and font standardisation (Arial)", 0),
    ("Document structure and headings", 0),
    ("Formats supported", -1),
    ("Word (.docx), Excel (.xlsx), PowerPoint (.pptx), Visio (.vsdx)", 0),
    ("Rules are data-driven", -1),
    ("Held in a SharePoint list and maintained by the business — no deployment to change a rule.", 0),
], size=17, gap=7)
footer(s)

# ================================================================ 7. ARCHITECTURE / TRUST BOUNDARY (diagram)
s = prs.slides.add_slide(BLANK)
title_bar(s, "How it fits together — and where your data lives", "Architecture & trust boundary")
# Mace tenant container
container(s, Inches(0.45), Inches(1.45), Inches(7.35), Inches(4.45), "MACE 365 TENANT  ·  your data, your controls", GREEN)
node(s, Inches(0.75), Inches(1.95), Inches(6.75), Inches(1.0),
     "Mace SharePoint  ·  Mace Way Control Centre", "Documents · Style rules · Validation results · Reports", fill=GREEN, tsz=13, ssz=10)
node(s, Inches(0.75), Inches(3.05), Inches(6.75), Inches(0.9),
     "Power Automate flow", "Runs as a Mace user · performs all SharePoint writes", fill=GREY, tsz=13, ssz=10)
node(s, Inches(0.75), Inches(4.05), Inches(6.75), Inches(0.9),
     "Microsoft Entra app registration", "Sites.Selected — one site only · admin-consented · revocable", fill=DARK, tsz=13, ssz=10)
# External Azure container
container(s, Inches(8.55), Inches(1.45), Inches(4.35), Inches(4.45), "EXTERNAL AZURE SUBSCRIPTION", DARK)
node(s, Inches(8.85), Inches(3.0), Inches(3.75), Inches(1.45),
     "Azure Function", "Validation engine\nStateless · in-memory · stores nothing", fill=DARK, tsz=14, ssz=11)
# cross-boundary arrows
arrow(s, MSO_SHAPE.RIGHT_ARROW, Inches(7.9), Inches(3.25), Inches(0.6), Inches(0.32), color=MUTED)
arrow(s, MSO_SHAPE.LEFT_ARROW, Inches(7.9), Inches(3.85), Inches(0.6), Inches(0.32), color=MUTED)
note_strip(s, Inches(0.45), Inches(6.05), Inches(12.45), Inches(0.8),
           "Only transient document content crosses the boundary — encrypted (TLS) and authenticated with a Mace-held key. "
           "Documents, rules, results and reports never leave the Mace tenant.", border=GREEN)
footer(s)

# ================================================================ 8. SECURITY HEADLINE
s = prs.slides.add_slide(BLANK)
title_bar(s, "Is the data secure? Yes.", "Security — the headline")
tf = textbox(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(1.0))
r = tf.paragraphs[0].add_run()
r.text = ("The validation engine runs in a separate Azure subscription — but your "
          "data never leaves the Mace tenant.")
_set(r, 20, TEXT, bold=True)
w, g, y0 = Inches(3.95), Inches(0.25), Inches(2.85)
card(s, Inches(0.7), y0, w, Inches(2.4), "Data stays in Mace",
     "Documents, rules, results and reports all live in Mace SharePoint. The engine never copies them out.", GREEN)
card(s, Inches(0.7) + (w + g), y0, w, Inches(2.4), "Least-privilege access",
     "Site-scoped (Sites.Selected), granted and revocable by Mace IT. It can see one site — nothing else.", GREY)
card(s, Inches(0.7) + 2 * (w + g), y0, w, Inches(2.4), "Nothing stored outside",
     "The engine is stateless: it processes in memory and discards. No database, no retention.", DARK)
footer(s)

# ================================================================ 9. SECURITY: RESIDENCY
s = prs.slides.add_slide(BLANK)
title_bar(s, "Your data stays in your tenant", "Security — data residency")
bullets(s, [
    "All documents, style rules, validation results and reports live in Mace SharePoint (Mace Way Control Centre). The engine never persists them elsewhere.",
    "The engine is stateless — documents are processed in memory and discarded. There is no database and no data retention in the external subscription.",
    "The Power Automate flow runs inside Mace, as a Mace user — every write to SharePoint happens within the Mace tenant.",
    "All traffic is encrypted in transit (TLS / HTTPS) end to end.",
    "External AI is currently disabled — no document content is sent to any third-party AI service.",
], gap=12)
footer(s)

# ================================================================ 10. SECURITY: ACCESS
s = prs.slides.add_slide(BLANK)
title_bar(s, "Access is least-privilege and Mace-controlled", "Security — access & control")
bullets(s, [
    "Access is granted through an Azure AD app registration on the Mace tenant, admin-consented by Mace IT.",
    "It uses Sites.Selected — scoped to a single site (Mace Way Control Centre), not tenant-wide. It cannot reach any other Mace data.",
    "The engine's HTTP endpoint is protected by function keys — only the Mace flow can invoke it.",
    "Mace can revoke access instantly — withdraw the site grant, disable the app, or rotate the key — with no dependency on the external subscription.",
    "Every request is logged with an audit trail and access-control checks (aligned to SOC 2 control families).",
], gap=11)
footer(s)

# ================================================================ 11. WHY HERE
s = prs.slides.add_slide(BLANK)
title_bar(s, 'Addressing "it\'s in your tenant"', "Security — hosting & path to production")
bullets(s, [
    "Today the engine runs in a personal Azure subscription. This was a deliberate pilot choice: rapid iteration with zero dependency on Mace infrastructure or change-control.",
    "The security model above holds regardless of where the engine runs — because the data and the controls live in Mace.",
    "The engine is tenant-agnostic. Productionising means redeploying the same code into a Mace-owned Azure subscription — no rewrite, no change to the data model.",
    "Recommendation: adopt the engine into Mace Azure under CDIO governance, with a nominated DTS owner.",
], gap=12)
footer(s)

# ================================================================ 12. PILOT -> PRODUCTION (diagram)
s = prs.slides.add_slide(BLANK)
title_bar(s, "From pilot to production — same code", "Path to production")
container(s, Inches(0.5), Inches(1.6), Inches(5.55), Inches(4.5), "TODAY  ·  PILOT", GREY)
node(s, Inches(0.85), Inches(2.2), Inches(4.85), Inches(1.45),
     "Mace 365 tenant", "Data · SharePoint · Power Automate flow", fill=GREEN, tsz=14, ssz=11)
arrow(s, MSO_SHAPE.DOWN_ARROW, Inches(2.95), Inches(3.72), Inches(0.35), Inches(0.5), color=MUTED)
node(s, Inches(0.85), Inches(4.3), Inches(4.85), Inches(1.3),
     "Validation engine", "Personal Azure subscription", fill=DARK, tsz=14, ssz=11)
# middle arrow
arrow(s, MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(3.55), Inches(0.95), Inches(0.5), color=GREY)
tb = textbox(s, Inches(5.95), Inches(3.05), Inches(1.45), Inches(0.5))
tb.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tb.paragraphs[0].add_run(); r.text = "Redeploy"; _set(r, 11, GREY, bold=True)
# target
container(s, Inches(7.25), Inches(1.6), Inches(5.55), Inches(4.5), "TARGET  ·  PRODUCTION", GREEN)
node(s, Inches(7.6), Inches(2.2), Inches(4.85), Inches(1.45),
     "Mace 365 tenant", "Data · SharePoint · Power Automate flow", fill=GREEN, tsz=14, ssz=11)
arrow(s, MSO_SHAPE.DOWN_ARROW, Inches(9.7), Inches(3.72), Inches(0.35), Inches(0.5), color=MUTED)
node(s, Inches(7.6), Inches(4.3), Inches(4.85), Inches(1.3),
     "Validation engine", "Mace Azure subscription", fill=GREEN, tsz=14, ssz=11)
note_strip(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.65),
           "Data residency and the security model are identical in both states — only the engine's hosting moves.", border=GREEN)
footer(s)

# ================================================================ 13. DEPLOYMENT
s = prs.slides.add_slide(BLANK)
title_bar(s, "Deployment & operations", "Deployment guide")
bullets(s, [
    ("Components", -1),
    ("Azure Function (Python, serverless / Flex Consumption) · Power Automate flow · SharePoint lists & libraries.", 0),
    ("Repeatable & governed", -1),
    ("Infrastructure-as-code (ARM template) and a CI/CD pipeline are defined; deployment is scripted, not manual.", 0),
    ("Configurable without redeploying", -1),
    ("Style rules and results are managed in SharePoint — the business changes rules without a release.", 0),
    ("Current state", -1),
    ("Piloted on Mace Way Control Centre and validated end-to-end on production documents.", 0),
], size=16, gap=6)
footer(s)

# ================================================================ 14. BUSINESS VALUE
s = prs.slides.add_slide(BLANK)
title_bar(s, "Why it matters", "Business value")
w, g, x0 = Inches(3.95), Inches(0.25), Inches(0.7)
vals = [
    ("Consistency & quality", "Every document held to the same standard, automatically.", GREY, Inches(1.6)),
    ("Speed", "Seconds, not manual review cycles — at the point of authoring.", GREEN, Inches(1.6)),
    ("Governance", "Auditable and aligned to the Control Centre's ISO 19650 ways of working.", GREY, Inches(1.6)),
    ("Low cost & scalable", "Serverless, pay-per-use; rules maintained by the business.", DARK, Inches(4.1)),
    ("A reusable pattern", "A foundation for automated document governance across Mace.", GREY, Inches(4.1)),
]
for i, (hd, bd, ac, yy) in enumerate(vals):
    card(s, x0 + (i % 3) * (w + g), yy, w, Inches(2.3), hd, bd, ac)
footer(s)

# ================================================================ 15. ROADMAP
s = prs.slides.add_slide(BLANK)
title_bar(s, "Roadmap", "What's next")
bullets(s, [
    ("Now", -1),
    ("Live pilot on Mace Way Control Centre — validation, auto-fix, branded reporting and audit, all working.", 0),
    ("Next", -1),
    ("Migrate the engine into Mace Azure; tighten the site grant to least-privilege (read-only); expand rule coverage.", 0),
    ("Later", -1),
    ("Optional AI-assisted language review with data-classification controls; roll out to further libraries and sites.", 0),
], size=18, gap=8)
footer(s)

# ================================================================ 16. THE ASK
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
bg.fill.solid(); bg.fill.fore_color.rgb = GREY; bg.line.fill.background(); bg.shadow.inherit = False
add_logo(s, Inches(0.9), Inches(0.7), Inches(0.95))
tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5))
r = tf.paragraphs[0].add_run(); r.text = "What I'm asking for"; _set(r, 34, WHITE, bold=True)
for t in [
    "Endorsement to productionise the engine within Mace Azure, under CDIO governance.",
    "A nominated owner in DTS / IT for the app registration and hosting.",
    "Support to grow this from a pilot into a Control-Centre-wide capability.",
]:
    p = tf.add_paragraph(); p.space_before = Pt(14); r = p.add_run(); r.text = "•  " + t; _set(r, 19, WHITE)
p = tf.add_paragraph(); p.space_before = Pt(26); r = p.add_run()
r.text = "Stephen Cummins  ·  stephen.cummins@macegroup.com"; _set(r, 14, RGBColor(0xDD, 0xDD, 0xDD))

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
